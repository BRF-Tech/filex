#!/usr/bin/env python3
"""
Generate the 12 demo fixtures expected by seed-example-fixtures.sh.

Idempotent — overwrites whatever is already at OUT_DIR. Used both for the
committed fixtures under e2e/fixtures/file-types/ and as the runtime
generator invoked by scripts/seed-example-fixtures.sh on the main host.

OUT_DIR may be overridden via the FIXTURE_OUT env var (default:
<repo>/e2e/fixtures/file-types).
"""

import json
import os
import struct
import sys
import zipfile
from pathlib import Path


def _default_out_dir() -> Path:
    here = Path(__file__).resolve().parent
    return here.parent / "e2e" / "fixtures" / "file-types"


OUT = Path(os.environ.get("FIXTURE_OUT", _default_out_dir()))
OUT.mkdir(parents=True, exist_ok=True)


# ---- Office (OOXML / ODF, all zip-based) -----------------------------------

def write_xlsx(p: Path) -> None:
    with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
<Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
</Types>""")
        z.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>""")
        z.writestr("xl/_rels/workbook.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
</Relationships>""")
        z.writestr("xl/workbook.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
<sheets><sheet name="Report" sheetId="1" r:id="rId1"/></sheets>
</workbook>""")
        z.writestr("xl/worksheets/sheet1.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
<sheetData><row r="1"><c r="A1" t="inlineStr"><is><t>Demo Report</t></is></c></row></sheetData>
</worksheet>""")


def write_docx(p: Path) -> None:
    with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""")
        z.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""")
        z.writestr("word/document.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
<w:body><w:p><w:r><w:t>Demo letter for FileManager preview routing.</w:t></w:r></w:p></w:body>
</w:document>""")


def write_pptx(p: Path) -> None:
    """Generate a PowerPoint fixture.

    The first attempt uses python-pptx, which produces a fully-formed
    archive (slideLayouts/, slideMasters/, theme/, app/core props…) —
    the kind of pptx LibreOffice Impress's writer can re-export to PDF
    via `soffice --convert-to pdf`. We use this path in production
    seeding (`scripts/seed-example-fixtures.sh` runs on `main`, where
    `pip install python-pptx` is already done) so the thumb pipeline
    actually has a working fixture.

    The fallback below is a hand-rolled minimal pptx — kept around so
    a developer running `python3 _gen_fixtures.py` on a vanilla machine
    without python-pptx installed still gets some bytes on disk. The
    frontend viewer tests open the fixture as an opaque blob and don't
    care whether it round-trips through Impress; only the thumbnail
    pipeline needs the richer form.
    """
    try:
        from pptx import Presentation  # type: ignore[import-not-found]

        prs = Presentation()
        title_layout = prs.slide_layouts[0]  # "Title Slide"
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = "Demo Presentation"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Routing fixture for filex preview viewers."
        prs.save(str(p))
        return
    except ImportError:
        pass

    # ---- Fallback: hand-rolled minimal pptx (no python-pptx on PATH) ----
    # NB: LibreOffice's Impress writer rejects pptx slides whose spTree
    # is empty — `soffice --convert-to pdf` exits 0 but emits no PDF
    # because there's nothing to lay out. We add a single text shape so
    # the converter has at least one drawable object, plus the
    # `<p:sldSz>` + `<p:notesSz>` size hints presentation.xml needs for
    # the slide to render at the standard 16:9 aspect. This still
    # doesn't satisfy Impress (no slideLayout/slideMaster references),
    # so on a python-pptx-less host the slides.pptx thumb stays in
    # state="failed" — re-seed from `main` to pick up the proper form.
    with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
<Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
</Types>""")
        z.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>""")
        z.writestr("ppt/_rels/presentation.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
</Relationships>""")
        z.writestr("ppt/presentation.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
<p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>
<p:sldSz cx="9144000" cy="6858000" type="screen4x3"/>
<p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>""")
        z.writestr("ppt/slides/slide1.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Title"/>
          <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
          <p:nvPr/>
        </p:nvSpPr>
        <p:spPr>
          <a:xfrm>
            <a:off x="685800" y="685800"/>
            <a:ext cx="7772400" cy="1143000"/>
          </a:xfrm>
          <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        </p:spPr>
        <p:txBody>
          <a:bodyPr anchor="ctr"/>
          <a:lstStyle/>
          <a:p>
            <a:pPr algn="ctr"/>
            <a:r>
              <a:rPr lang="en-US" sz="4400" b="1"/>
              <a:t>Demo Presentation</a:t>
            </a:r>
          </a:p>
        </p:txBody>
      </p:sp>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="3" name="Body"/>
          <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
          <p:nvPr/>
        </p:nvSpPr>
        <p:spPr>
          <a:xfrm>
            <a:off x="685800" y="2286000"/>
            <a:ext cx="7772400" cy="3429000"/>
          </a:xfrm>
          <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        </p:spPr>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          <a:p>
            <a:r>
              <a:rPr lang="en-US" sz="2400"/>
              <a:t>Routing fixture for filex preview viewers.</a:t>
            </a:r>
          </a:p>
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
</p:sld>""")


def _odf_zip(p: Path, mimetype: str, content_xml_body: str) -> None:
    """Write an ODF package — `mimetype` MUST be the first entry, STORED, no extra."""
    with zipfile.ZipFile(p, "w") as z:
        info = zipfile.ZipInfo("mimetype")
        info.compress_type = zipfile.ZIP_STORED
        z.writestr(info, mimetype)
        z.writestr("META-INF/manifest.xml", f"""<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">
<manifest:file-entry manifest:full-path="/" manifest:media-type="{mimetype}"/>
<manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>
</manifest:manifest>""")
        z.writestr("content.xml", content_xml_body)


def write_odt(p: Path) -> None:
    body = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">
<office:body><office:text><text:p>Demo notes - LibreOffice text doc for preview routing.</text:p></office:text></office:body>
</office:document-content>"""
    _odf_zip(p, "application/vnd.oasis.opendocument.text", body)


def write_ods(p: Path) -> None:
    body = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0" xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">
<office:body><office:spreadsheet>
<table:table table:name="Budget">
<table:table-row><table:table-cell office:value-type="string"><text:p>Budget</text:p></table:table-cell></table:table-row>
</table:table></office:spreadsheet></office:body>
</office:document-content>"""
    _odf_zip(p, "application/vnd.oasis.opendocument.spreadsheet", body)


# ---- Diagram / notebook -----------------------------------------------------

def write_drawio(p: Path) -> None:
    p.write_text("""<mxfile host="app.diagrams.net" modified="2026-05-07T00:00:00.000Z" agent="filemanager-fixture" version="22.0.0">
  <diagram name="Page-1" id="demo">
    <mxGraphModel dx="800" dy="600" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="850" pageHeight="1100" math="0" shadow="0">
      <root>
        <mxCell id="0" />
        <mxCell id="1" parent="0" />
        <mxCell id="2" value="Hello" style="rounded=0;whiteSpace=wrap;html=1;" vertex="1" parent="1">
          <mxGeometry x="160" y="120" width="120" height="60" as="geometry" />
        </mxCell>
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>
""", encoding="utf-8")


def write_mmd(p: Path) -> None:
    p.write_text("""graph TD
    A[Start] --> B{Decision}
    B -->|yes| C[Do thing]
    B -->|no| D[Skip]
    C --> E[End]
    D --> E
""", encoding="utf-8")


def write_ipynb(p: Path) -> None:
    payload = {
        "cells": [
            {
                "cell_type": "markdown",
                "metadata": {},
                "source": ["# Demo notebook\n", "Routing fixture for IpynbViewer."],
            },
            {
                "cell_type": "code",
                "execution_count": 1,
                "metadata": {},
                "outputs": [
                    {"output_type": "stream", "name": "stdout", "text": ["hello\n"]}
                ],
                "source": ["print('hello')"],
            },
        ],
        "metadata": {
            "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"}
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    p.write_text(json.dumps(payload, indent=2), encoding="utf-8")


# ---- 3D --------------------------------------------------------------------

def write_stl(p: Path) -> None:
    p.write_text("""solid cube
  facet normal 0 0 1
    outer loop
      vertex 0 0 0
      vertex 1 0 0
      vertex 0 1 0
    endloop
  endfacet
endsolid cube
""", encoding="utf-8")


def write_obj(p: Path) -> None:
    p.write_text("""# minimal OBJ
v 0.0 0.0 0.0
v 1.0 0.0 0.0
v 0.0 1.0 0.0
f 1 2 3
""", encoding="utf-8")


def write_glb(p: Path) -> None:
    """Minimal valid GLB (glTF binary): header + JSON chunk only."""
    json_body = b'{"asset":{"version":"2.0"},"scenes":[{"nodes":[]}],"scene":0,"nodes":[],"meshes":[]}'
    pad = (4 - len(json_body) % 4) % 4
    json_body += b" " * pad
    json_chunk = struct.pack("<I", len(json_body)) + b"JSON" + json_body
    total = 12 + len(json_chunk)
    header = b"glTF" + struct.pack("<II", 2, total)
    p.write_bytes(header + json_chunk)


# ---- EPUB ------------------------------------------------------------------

def write_epub(p: Path) -> None:
    with zipfile.ZipFile(p, "w") as z:
        info = zipfile.ZipInfo("mimetype")
        info.compress_type = zipfile.ZIP_STORED
        z.writestr(info, "application/epub+zip")
        z.writestr("META-INF/container.xml", """<?xml version="1.0" encoding="UTF-8"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>""")
        z.writestr("OEBPS/content.opf", """<?xml version="1.0" encoding="UTF-8"?>
<package version="3.0" xmlns="http://www.idpf.org/2007/opf" unique-identifier="bookid">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:identifier id="bookid">demo-book-001</dc:identifier>
    <dc:title>Demo Book</dc:title>
    <dc:language>en</dc:language>
    <meta property="dcterms:modified">2026-05-07T00:00:00Z</meta>
  </metadata>
  <manifest>
    <item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/>
    <item id="ch1" href="ch1.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine>
    <itemref idref="ch1"/>
  </spine>
</package>""")
        z.writestr("OEBPS/nav.xhtml", """<?xml version="1.0" encoding="UTF-8"?>
<html xmlns="http://www.w3.org/1999/xhtml" xmlns:epub="http://www.idpf.org/2007/ops">
<head><title>Nav</title></head>
<body><nav epub:type="toc"><ol><li><a href="ch1.xhtml">Chapter 1</a></li></ol></nav></body>
</html>""")
        z.writestr("OEBPS/ch1.xhtml", """<?xml version="1.0" encoding="UTF-8"?>
<html xmlns="http://www.w3.org/1999/xhtml">
<head><title>Chapter 1</title></head>
<body><h1>Demo Book</h1><p>Routing fixture for EpubViewer.</p></body>
</html>""")


# ---- PSD (1x1 RGB, valid 8BPS header) -------------------------------------

def write_psd(p: Path) -> None:
    sig = b"8BPS"
    version = struct.pack(">H", 1)
    reserved = b"\x00" * 6
    channels = struct.pack(">H", 3)
    height = struct.pack(">I", 1)
    width = struct.pack(">I", 1)
    depth = struct.pack(">H", 8)
    color_mode = struct.pack(">H", 3)  # RGB
    header = sig + version + reserved + channels + height + width + depth + color_mode
    color_mode_data = struct.pack(">I", 0)
    image_resources = struct.pack(">I", 0)
    layer_info = struct.pack(">I", 0)
    image_data = struct.pack(">H", 0) + b"\xff\x80\x40"
    p.write_bytes(header + color_mode_data + image_resources + layer_info + image_data)


# ---- TIFF (1x1 grayscale, hand-crafted little-endian) ---------------------

def write_tiff(p: Path) -> None:
    out = bytearray()
    out += b"II"
    out += struct.pack("<H", 42)
    out += struct.pack("<I", 8)
    entries = []
    entries.append(struct.pack("<HHI", 256, 3, 1) + struct.pack("<HH", 1, 0))   # ImageWidth=1
    entries.append(struct.pack("<HHI", 257, 3, 1) + struct.pack("<HH", 1, 0))   # ImageLength=1
    entries.append(struct.pack("<HHI", 258, 3, 1) + struct.pack("<HH", 8, 0))   # BitsPerSample=8
    entries.append(struct.pack("<HHI", 259, 3, 1) + struct.pack("<HH", 1, 0))   # Compression=none
    entries.append(struct.pack("<HHI", 262, 3, 1) + struct.pack("<HH", 1, 0))   # Photometric=BlackIsZero
    so_idx = len(entries)
    entries.append(struct.pack("<HHI", 273, 4, 1) + struct.pack("<I", 0))       # StripOffsets — patch later
    entries.append(struct.pack("<HHI", 277, 3, 1) + struct.pack("<HH", 1, 0))   # SamplesPerPixel=1
    entries.append(struct.pack("<HHI", 278, 3, 1) + struct.pack("<HH", 1, 0))   # RowsPerStrip=1
    entries.append(struct.pack("<HHI", 279, 4, 1) + struct.pack("<I", 1))       # StripByteCounts=1
    n = len(entries)
    ifd_size = 2 + n * 12 + 4
    strip_offset = 8 + ifd_size
    entries[so_idx] = struct.pack("<HHI", 273, 4, 1) + struct.pack("<I", strip_offset)
    out += struct.pack("<H", n)
    for e in entries:
        out += e
    out += struct.pack("<I", 0)  # next-IFD offset
    out += b"\x80"               # one grayscale pixel
    p.write_bytes(bytes(out))


# ---- Driver ---------------------------------------------------------------

WRITERS = [
    ("report.xlsx", write_xlsx),
    ("letter.docx", write_docx),
    ("slides.pptx", write_pptx),
    ("notes.odt", write_odt),
    ("budget.ods", write_ods),
    ("diagram.drawio", write_drawio),
    ("flow.mmd", write_mmd),
    ("notebook.ipynb", write_ipynb),
    ("cube.stl", write_stl),
    ("cube.obj", write_obj),
    ("cube.glb", write_glb),
    ("book.epub", write_epub),
    ("layered.psd", write_psd),
    ("scan.tiff", write_tiff),
]


def main() -> int:
    for name, fn in WRITERS:
        path = OUT / name
        fn(path)
    print(f"[fixtures] wrote {len(WRITERS)} files to {OUT}")
    for name, _ in WRITERS:
        path = OUT / name
        print(f"  {path.stat().st_size:>8}  {name}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
